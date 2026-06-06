"""Generate the techbold PowerPoint template deck.

This helper creates pitch-deck slide layouts. It is not part of the runtime app.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

# Brand tokens (from DESIGN.md)
NAVY      = RGBColor(0x26, 0x2b, 0x4b)   # #262b4b  structural navy
NAVY_MID  = RGBColor(0x37, 0x36, 0x52)   # #373652  interactive navy
GOLD      = RGBColor(0xfc, 0xb5, 0x14)   # #fcb514  techbold gold
GOLD_LT   = RGBColor(0xf7, 0xcf, 0x42)   # #f7cf42  gold light
WARM_BG   = RGBColor(0xf3, 0xf2, 0xea)   # #f3f2ea  warm surface
WHITE     = RGBColor(0xff, 0xff, 0xff)
OFF_WHITE = RGBColor(0xfb, 0xfb, 0xfb)
BORDER    = RGBColor(0xe5, 0xe7, 0xea)
MUTED     = RGBColor(0x80, 0x80, 0x80)

# Font stacks — PowerPoint falls back to system fonts, so we map to closest
# available equivalents. Barlow Condensed and Inter ship on many systems;
# if absent, Arial Narrow / Arial are the safe fallbacks.
FONT_DISPLAY = "Barlow Condensed"   # ≈ ProximaNova Condensed Black
FONT_BODY    = "Inter"              # ≈ Denim WD
FONT_MONO    = "JetBrains Mono"     # for code/command slides

# Slide dimensions: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

MARGIN = Inches(0.6)


# Helpers

def rgb_hex(r: RGBColor) -> str:
    """Convert a PowerPoint RGB color to a hex string."""

    return f"{r[0]:02X}{r[1]:02X}{r[2]:02X}"


def solid_fill(shape, color: RGBColor):
    """Fill shape background with a solid color."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def no_fill(shape):
    """Remove the shape fill."""

    shape.fill.background()


def no_line(shape):
    """Remove the shape outline."""

    shape.line.fill.background()


def add_rect(slide, left, top, width, height, color: RGBColor):
    """Add a solid rectangle to a slide."""

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    solid_fill(shape, color)
    no_line(shape)
    return shape


def add_textbox(slide, left, top, width, height):
    """Add a textbox with the given PowerPoint coordinates."""

    return slide.shapes.add_textbox(left, top, width, height)


def set_para(para, align=PP_ALIGN.LEFT, space_before=0, space_after=0):
    """Apply common paragraph spacing and alignment."""

    para.alignment = align
    para.space_before = Pt(space_before)
    para.space_after = Pt(space_after)


def add_run(para, text, font_name, size_pt, bold=False, color=None, italic=False, caps=False):
    """Add one styled text run to a paragraph."""

    run = para.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    if caps:
        run.font.all_caps = True
    return run


def label(slide, left, top, width, text, color=MUTED):
    """Small uppercase tracking label."""
    tb = add_textbox(slide, left, top, width, Inches(0.25))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    set_para(p)
    add_run(p, text.upper(), FONT_BODY, 7, bold=True, color=color, caps=True)
    return tb


def gold_bar(slide, top=Inches(0), height=Inches(0.06)):
    """Thin horizontal gold accent bar."""
    return add_rect(slide, 0, top, W, height, GOLD)


def navy_stripe(slide, left=Inches(0), width=Inches(0.08)):
    """Thin vertical navy accent on left edge."""
    return add_rect(slide, left, 0, width, H, NAVY)


# Slide builders

def build_title_slide(prs):
    """
    Full-bleed navy slide. Gold bar top. White wordmark area.
    Layout: large display title, subtitle, gold CTA label bottom-right.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    add_rect(slide, 0, 0, W, H, NAVY)

    # Gold accent bar — top
    gold_bar(slide, top=0, height=Inches(0.07))

    # Gold accent bar — bottom
    gold_bar(slide, top=H - Inches(0.07), height=Inches(0.07))

    # Diagonal gold block — decorative right panel
    add_rect(slide, W - Inches(3.8), 0, Inches(3.8), H, NAVY_MID)

    # Gold vertical stripe separating panels
    add_rect(slide, W - Inches(3.85), 0, Inches(0.06), H, GOLD)

    # Title area
    tb = add_textbox(slide, MARGIN, Inches(2.2), Inches(8.5), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    set_para(p, align=PP_ALIGN.LEFT, space_after=6)
    add_run(p, "Presentation Title", FONT_DISPLAY, 48, bold=True, color=WHITE, caps=True)

    p2 = tf.add_paragraph()
    set_para(p2, align=PP_ALIGN.LEFT)
    add_run(p2, "Subtitle or project name goes here", FONT_BODY, 18, color=RGBColor(0xb0, 0xb8, 0xd4))

    # Metadata block bottom-left
    meta_tb = add_textbox(slide, MARGIN, H - Inches(1.4), Inches(5), Inches(0.9))
    meta_tf = meta_tb.text_frame
    p3 = meta_tf.paragraphs[0]
    set_para(p3)
    add_run(p3, "DATE  ·  AUTHOR  ·  VERSION", FONT_BODY, 8, bold=True,
            color=RGBColor(0x70, 0x78, 0x9a), caps=True)

    # Logo placeholder — right panel
    logo_tb = add_textbox(slide, W - Inches(3.4), Inches(3.1), Inches(2.8), Inches(0.7))
    logo_tf = logo_tb.text_frame
    p4 = logo_tf.paragraphs[0]
    set_para(p4, align=PP_ALIGN.CENTER)
    add_run(p4, "techbold", FONT_DISPLAY, 32, bold=True, color=WHITE, caps=True)

    # Tagline below logo
    tag_tb = add_textbox(slide, W - Inches(3.4), Inches(3.9), Inches(2.8), Inches(0.4))
    tag_tf = tag_tb.text_frame
    p5 = tag_tf.paragraphs[0]
    set_para(p5, align=PP_ALIGN.CENTER)
    add_run(p5, "die Extrameile ist unser Standard", FONT_BODY, 9,
            color=RGBColor(0x90, 0x98, 0xb8), italic=True)

    return slide


def build_section_divider(prs):
    """
    Gold-dominant section break. Navy text on gold. Bold.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gold background
    add_rect(slide, 0, 0, W, H, GOLD)

    # Navy block left strip
    add_rect(slide, 0, 0, Inches(1.2), H, NAVY)

    # Section number (vertical, rotated feel via wide box)
    num_tb = add_textbox(slide, Inches(0.15), Inches(2.8), Inches(0.9), Inches(1.2))
    num_tf = num_tb.text_frame
    p = num_tf.paragraphs[0]
    set_para(p, align=PP_ALIGN.CENTER)
    add_run(p, "01", FONT_DISPLAY, 36, bold=True, color=GOLD)

    # Section label
    lbl_tb = add_textbox(slide, Inches(1.6), Inches(1.8), Inches(9), Inches(0.5))
    lbl_tf = lbl_tb.text_frame
    p2 = lbl_tf.paragraphs[0]
    set_para(p2)
    add_run(p2, "SECTION TITLE", FONT_BODY, 10, bold=True, color=NAVY, caps=True)

    # Large section heading
    head_tb = add_textbox(slide, Inches(1.6), Inches(2.4), Inches(9), Inches(2.5))
    head_tf = head_tb.text_frame
    head_tf.word_wrap = True
    p3 = head_tf.paragraphs[0]
    set_para(p3, space_after=8)
    add_run(p3, "Section Heading\nGoes Here", FONT_DISPLAY, 52, bold=True, color=NAVY, caps=True)

    # Bottom navy bar
    add_rect(slide, 0, H - Inches(0.07), W, Inches(0.07), NAVY)

    return slide


def build_content_slide(prs):
    """
    Standard content: warm-surface bg, navy header bar, body area.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Warm surface background
    add_rect(slide, 0, 0, W, H, WARM_BG)

    # Navy header bar
    add_rect(slide, 0, 0, W, Inches(1.15), NAVY)

    # Gold accent stripe below header
    gold_bar(slide, top=Inches(1.15), height=Inches(0.05))

    # Slide title in header
    title_tb = add_textbox(slide, MARGIN, Inches(0.18), Inches(10), Inches(0.75))
    title_tf = title_tb.text_frame
    p = title_tf.paragraphs[0]
    set_para(p, align=PP_ALIGN.LEFT)
    add_run(p, "Slide Title", FONT_DISPLAY, 30, bold=True, color=WHITE, caps=True)

    # Section breadcrumb top-right
    bc_tb = add_textbox(slide, W - Inches(3), Inches(0.38), Inches(2.4), Inches(0.35))
    bc_tf = bc_tb.text_frame
    p2 = bc_tf.paragraphs[0]
    set_para(p2, align=PP_ALIGN.RIGHT)
    add_run(p2, "SECTION NAME", FONT_BODY, 7, bold=True,
            color=RGBColor(0x70, 0x78, 0x9a), caps=True)

    # Body text area
    body_tb = add_textbox(slide, MARGIN, Inches(1.45), Inches(12.1), Inches(5.6))
    body_tf = body_tb.text_frame
    body_tf.word_wrap = True

    p3 = body_tf.paragraphs[0]
    set_para(p3, space_after=10)
    add_run(p3, "Body text goes here. Keep lines to 65–72 characters for readability. "
                "Body is Inter Regular 16pt on warm-surface background.", FONT_BODY, 16, color=NAVY)

    p4 = body_tf.add_paragraph()
    set_para(p4, space_before=4, space_after=4)
    add_run(p4, "• Bullet point one with supporting detail", FONT_BODY, 15, color=NAVY_MID)

    p5 = body_tf.add_paragraph()
    set_para(p5, space_before=4, space_after=4)
    add_run(p5, "• Bullet point two with supporting detail", FONT_BODY, 15, color=NAVY_MID)

    p6 = body_tf.add_paragraph()
    set_para(p6, space_before=4, space_after=4)
    add_run(p6, "• Bullet point three with supporting detail", FONT_BODY, 15, color=NAVY_MID)

    # Slide number — bottom right
    num_tb = add_textbox(slide, W - Inches(1), H - Inches(0.4), Inches(0.8), Inches(0.3))
    num_tf = num_tb.text_frame
    p7 = num_tf.paragraphs[0]
    set_para(p7, align=PP_ALIGN.RIGHT)
    add_run(p7, "01", FONT_BODY, 9, color=MUTED)

    return slide


def build_two_col_slide(prs):
    """
    Two-column layout. Left: navy card. Right: warm content.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, W, H, WARM_BG)

    # Header
    add_rect(slide, 0, 0, W, Inches(1.15), NAVY)
    gold_bar(slide, top=Inches(1.15), height=Inches(0.05))

    title_tb = add_textbox(slide, MARGIN, Inches(0.18), Inches(10), Inches(0.75))
    title_tf = title_tb.text_frame
    p = title_tf.paragraphs[0]
    set_para(p)
    add_run(p, "Two-Column Layout", FONT_DISPLAY, 30, bold=True, color=WHITE, caps=True)

    COL_GAP = Inches(0.25)
    COL_TOP = Inches(1.35)
    COL_H = H - COL_TOP - Inches(0.3)
    COL_W = (W - MARGIN * 2 - COL_GAP) / 2

    # Left navy card
    left_card = add_rect(slide, MARGIN, COL_TOP, COL_W, COL_H, NAVY)
    # Gold top accent on card
    add_rect(slide, MARGIN, COL_TOP, COL_W, Inches(0.055), GOLD)

    l_tb = add_textbox(slide, MARGIN + Inches(0.3), COL_TOP + Inches(0.35),
                       COL_W - Inches(0.6), COL_H - Inches(0.5))
    l_tf = l_tb.text_frame
    l_tf.word_wrap = True

    lp = l_tf.paragraphs[0]
    set_para(lp, space_after=12)
    add_run(lp, "Key Point", FONT_DISPLAY, 22, bold=True, color=GOLD, caps=True)

    lp2 = l_tf.add_paragraph()
    set_para(lp2)
    add_run(lp2, "Supporting text or metric. Keep it short and impactful. "
                 "Use this column for the key claim.", FONT_BODY, 14, color=WHITE)

    # Right content area
    r_left = MARGIN + COL_W + COL_GAP
    r_tb = add_textbox(slide, r_left, COL_TOP + Inches(0.25),
                       COL_W, COL_H - Inches(0.35))
    r_tf = r_tb.text_frame
    r_tf.word_wrap = True

    rp = r_tf.paragraphs[0]
    set_para(rp, space_after=10)
    add_run(rp, "Detail heading", FONT_BODY, 16, bold=True, color=NAVY)

    for txt in [
        "Supporting detail one that elaborates on the key point.",
        "Supporting detail two with additional context or data.",
        "Supporting detail three closing the argument.",
    ]:
        rp_n = r_tf.add_paragraph()
        set_para(rp_n, space_before=6, space_after=2)
        add_run(rp_n, f"• {txt}", FONT_BODY, 13, color=NAVY_MID)

    # Slide number
    num_tb = add_textbox(slide, W - Inches(1), H - Inches(0.4), Inches(0.8), Inches(0.3))
    p_n = num_tb.text_frame.paragraphs[0]
    set_para(p_n, align=PP_ALIGN.RIGHT)
    add_run(p_n, "02", FONT_BODY, 9, color=MUTED)

    return slide


def build_metrics_slide(prs):
    """
    3-up stat cards. Navy bg. Gold numbers.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, W, H, NAVY)
    gold_bar(slide, top=0, height=Inches(0.07))

    title_tb = add_textbox(slide, MARGIN, Inches(0.3), Inches(11), Inches(0.75))
    p = title_tb.text_frame.paragraphs[0]
    set_para(p)
    add_run(p, "Key Metrics", FONT_DISPLAY, 34, bold=True, color=WHITE, caps=True)

    CARD_TOP = Inches(1.4)
    CARD_H = Inches(4.5)
    CARD_GAP = Inches(0.3)
    N = 3
    CARD_W = (W - MARGIN * 2 - CARD_GAP * (N - 1)) / N

    stats = [
        ("170+", "IT Specialists", "Dedicated technicians\nacross all service areas"),
        ("1,200", "Customers Served", "Spanning 27 industries\nand counting"),
        ("4.68 / 5", "Satisfaction Score", "Based on last 12 months\nof service feedback"),
    ]

    for i, (number, label_txt, body_txt) in enumerate(stats):
        left = MARGIN + i * (CARD_W + CARD_GAP)

        # Card background
        card = add_rect(slide, left, CARD_TOP, CARD_W, CARD_H, NAVY_MID)

        # Gold top bar on card
        add_rect(slide, left, CARD_TOP, CARD_W, Inches(0.055), GOLD)

        # Number
        n_tb = add_textbox(slide, left + Inches(0.3), CARD_TOP + Inches(0.4),
                           CARD_W - Inches(0.6), Inches(1.4))
        p_n = n_tb.text_frame.paragraphs[0]
        set_para(p_n)
        add_run(p_n, number, FONT_DISPLAY, 52, bold=True, color=GOLD)

        # Label
        l_tb = add_textbox(slide, left + Inches(0.3), CARD_TOP + Inches(1.85),
                           CARD_W - Inches(0.6), Inches(0.45))
        p_l = l_tb.text_frame.paragraphs[0]
        set_para(p_l)
        add_run(p_l, label_txt.upper(), FONT_BODY, 9, bold=True, color=GOLD_LT, caps=True)

        # Body
        b_tb = add_textbox(slide, left + Inches(0.3), CARD_TOP + Inches(2.45),
                           CARD_W - Inches(0.6), Inches(1.6))
        b_tf = b_tb.text_frame
        b_tf.word_wrap = True
        p_b = b_tf.paragraphs[0]
        set_para(p_b)
        add_run(p_b, body_txt, FONT_BODY, 12, color=RGBColor(0xb0, 0xb8, 0xd4))

    # Gold bottom bar
    gold_bar(slide, top=H - Inches(0.07), height=Inches(0.07))

    return slide


def build_quote_slide(prs):
    """
    Full navy slide with large pull quote. Gold opening mark.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, NAVY)

    # Gold left strip
    add_rect(slide, 0, 0, Inches(0.08), H, GOLD)

    # Gold quotation mark
    qm_tb = add_textbox(slide, Inches(1.2), Inches(1.2), Inches(2), Inches(2))
    p_qm = qm_tb.text_frame.paragraphs[0]
    set_para(p_qm)
    add_run(p_qm, "“", FONT_DISPLAY, 120, bold=True, color=GOLD)

    # Quote text
    qt_tb = add_textbox(slide, Inches(1.4), Inches(2.2), Inches(10.5), Inches(3))
    qt_tf = qt_tb.text_frame
    qt_tf.word_wrap = True
    p_qt = qt_tf.paragraphs[0]
    set_para(p_qt, space_after=16)
    add_run(p_qt, "Die Extrameile ist unser Standard.", FONT_DISPLAY, 40, bold=True, color=WHITE)

    p_attr = qt_tf.add_paragraph()
    set_para(p_attr)
    add_run(p_attr, "— techbold", FONT_BODY, 14, color=RGBColor(0x90, 0x98, 0xb8))

    gold_bar(slide, top=H - Inches(0.07), height=Inches(0.07))
    return slide


def build_agenda_slide(prs):
    """
    Agenda / table of contents. Numbered items, gold numbers.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, WARM_BG)

    # Left navy panel
    add_rect(slide, 0, 0, Inches(3.6), H, NAVY)
    add_rect(slide, Inches(3.6), 0, Inches(0.06), H, GOLD)

    # Panel title
    pt_tb = add_textbox(slide, Inches(0.4), Inches(2.5), Inches(2.8), Inches(2))
    pt_tf = pt_tb.text_frame
    pt_tf.word_wrap = True
    p_pt = pt_tf.paragraphs[0]
    set_para(p_pt)
    add_run(p_pt, "AGENDA", FONT_DISPLAY, 38, bold=True, color=WHITE, caps=True)

    p_pt2 = pt_tf.add_paragraph()
    set_para(p_pt2, space_before=8)
    add_run(p_pt2, "Today's session", FONT_BODY, 13, color=RGBColor(0x90, 0x98, 0xb8))

    # Agenda items
    items = [
        ("01", "Introduction & Context"),
        ("02", "Current State Analysis"),
        ("03", "Proposed Solution"),
        ("04", "Implementation Roadmap"),
        ("05", "Q & A"),
    ]

    ITEM_LEFT = Inches(4.1)
    ITEM_TOP = Inches(1.1)
    ITEM_H = Inches(1.0)

    for i, (num, text) in enumerate(items):
        top = ITEM_TOP + i * ITEM_H

        # Number
        n_tb = add_textbox(slide, ITEM_LEFT, top, Inches(0.8), Inches(0.75))
        p_n = n_tb.text_frame.paragraphs[0]
        set_para(p_n, align=PP_ALIGN.LEFT)
        add_run(p_n, num, FONT_DISPLAY, 28, bold=True, color=GOLD)

        # Item text
        t_tb = add_textbox(slide, ITEM_LEFT + Inches(0.9), top + Inches(0.08),
                           Inches(7.8), Inches(0.6))
        p_t = t_tb.text_frame.paragraphs[0]
        set_para(p_t, align=PP_ALIGN.LEFT)
        add_run(p_t, text, FONT_BODY, 18, bold=(i == 0),
                color=NAVY if i != 0 else NAVY_MID)

        # Divider line (except after last)
        if i < len(items) - 1:
            add_rect(slide, ITEM_LEFT, top + Inches(0.82),
                     Inches(8.6), Pt(0.75), BORDER)

    return slide


def build_closing_slide(prs):
    """
    End / thank-you slide. Full navy. Contact info.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, NAVY)

    gold_bar(slide, top=0, height=Inches(0.07))

    # Large brand word
    brand_tb = add_textbox(slide, MARGIN, Inches(1.8), Inches(8), Inches(2))
    brand_tf = brand_tb.text_frame
    brand_tf.word_wrap = False
    p_b = brand_tf.paragraphs[0]
    set_para(p_b)
    add_run(p_b, "techbold", FONT_DISPLAY, 80, bold=True, color=WHITE, caps=True)

    # Gold underline
    add_rect(slide, MARGIN, Inches(4.1), Inches(5.5), Inches(0.07), GOLD)

    # Tagline
    tag_tb = add_textbox(slide, MARGIN, Inches(4.35), Inches(8), Inches(0.5))
    p_tag = tag_tb.text_frame.paragraphs[0]
    set_para(p_tag)
    add_run(p_tag, "die Extrameile ist unser Standard", FONT_BODY, 16,
            italic=True, color=RGBColor(0x90, 0x98, 0xb8))

    # Contact block
    contact_tb = add_textbox(slide, MARGIN, Inches(5.3), Inches(6), Inches(1.2))
    contact_tf = contact_tb.text_frame
    contact_tf.word_wrap = False

    for line, size in [("www.techbold.at", 14), ("+43 59 555", 13)]:
        p_c = contact_tf.add_paragraph() if line != "www.techbold.at" else contact_tf.paragraphs[0]
        set_para(p_c, space_after=4)
        add_run(p_c, line, FONT_BODY, size, color=RGBColor(0xb0, 0xb8, 0xd4))

    gold_bar(slide, top=H - Inches(0.07), height=Inches(0.07))

    return slide


# Main

def main():
    """Build the template deck and save it to the configured output path."""

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    build_title_slide(prs)
    build_agenda_slide(prs)
    build_section_divider(prs)
    build_content_slide(prs)
    build_two_col_slide(prs)
    build_metrics_slide(prs)
    build_quote_slide(prs)
    build_closing_slide(prs)

    out = "/Users/julianwollner/Documents/GitHub/techbold_track_start_hack/techbold_template.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
